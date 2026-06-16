#!/usr/bin/env python3
"""从 timeline.json 生成讲解 PPT。

每句一页：原文大字 + 关键词卡片 + 译文。

用法：
  python3 script/build_pptx.py "诫子书 全文"
  python3 script/build_pptx.py "诫子书" --out 诫子书.pptx
"""

import argparse
import base64
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
TIMELINE_DIR = PROJECT_ROOT / "data" / "timelines"
sys.path.insert(0, str(PROJECT_ROOT))
import generate


def _split_course_chapter(lesson: str) -> tuple[str, str]:
    """从 lesson 名提取课程名和章节名（和 build_timeline.py 一致）。"""
    shishuo_chapters = {
        "德行篇", "言语篇", "政事篇", "文学篇", "方正篇", "雅量篇",
        "识鉴篇", "赏誉篇", "品藻篇", "规箴篇", "捷悟篇", "夙慧篇",
        "豪爽篇", "容止篇", "自新篇", "俭啬篇", "汰侈篇", "忿狷篇",
        "情礼篇", "黜免篇", "俭吝篇", "惑溺篇", "仇隙篇", "任诞篇",
        "伤逝篇", "栖逸篇", "贤媛篇", "术解篇", "巧艺篇", "知惧篇",
        "企羡篇",
    }
    for chapter in shishuo_chapters:
        if chapter in lesson:
            return ("世说新语精读", chapter)
    return (lesson, lesson)

# 配色
INK = RGBColor(0x3A, 0x2F, 0x1F)        # 墨色
WARM_CREAM = RGBColor(0xF5, 0xF0, 0xE8) # 暖米背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xC4, 0xA9, 0x7D)       # 金色
GRAY = RGBColor(0x8C, 0x7A, 0x6B)       # 灰色
DARK_GOLD = RGBColor(0x5C, 0x3D, 0x2E)   # 暗金色

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

TAG_COLORS = [
    RGBColor(0xFA, 0xF6, 0xED),  # 米白
    RGBColor(0xFB, 0xF0, 0xE0),  # 浅杏
    RGBColor(0xF5, 0xEB, 0xE6),  # 浅驼
]


def _add_text_box(slide, left, top, width, height, text,
                  font_size=14, bold=False, color=INK, align=PP_ALIGN.CENTER,
                  font_name="PingFang SC"):
    """添加文本框。"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_rounded_rect(slide, left, top, width, height, fill_color):
    """圆角矩形。"""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# 1x1 透明 PNG，作为音频媒体对象的占位缩略图（移到画面外不显示）
_ICON_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR42mNk"
    "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)

# 自动播放音频：沿用 python-pptx/PowerPoint 原生 video media timing 结构，
# 只把开始条件从 delay="indefinite" 改为 delay="0"。
_TIMING_TPL = (
    '<p:timing><p:tnLst><p:par>'
    '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
    '<p:childTnLst><p:video><p:cMediaNode vol="80000">'
    '<p:cTn id="2" fill="hold" display="0">'
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    '</p:cTn>'
    '<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
    '</p:cMediaNode></p:video></p:childTnLst>'
    '</p:cTn></p:par></p:tnLst></p:timing>'
)


def _embed_audio_post(out_path: Path, audio_map: dict[int, Path],
                     intro_duration_ms: int = 3000) -> None:
    """后处理 pptx：注入 PowerPoint 原生媒体对象，并设为自动播放。

    关键点：PowerPoint 导出视频只采集幻灯片上的自动播放媒体对象。
    这里按 python-pptx add_movie 生成的结构写入 MP3：video rel + 2007 media rel + poster image。

    Args:
        out_path: 输出 PPTX 路径
        audio_map: {slide_index(1-based): audio_file_path}
        intro_duration_ms: 封面页停留时长（毫秒），默认 3000
    """
    import zipfile
    import shutil
    import re
    from mutagen.mp3 import MP3

    tmp = out_path.with_suffix(".tmp")
    shutil.move(str(out_path), str(tmp))

    adv_times = {}  # slide_num → advTm in ms
    mp3_map = {}    # slide_num → mp3 bytes
    for sn, audio_path in audio_map.items():
        if audio_path.exists():
            mp3_map[sn] = audio_path.read_bytes()
            adv_times[sn] = int(MP3(audio_path).info.length * 1000)
    
    # 封面页停留时长（从外部传入，等于导入语音频时长）
    if 1 not in adv_times:
        adv_times[1] = intro_duration_ms  # 使用传入的时长

    slide_xmls = {}   # slide_num → xml_str
    slide_rels = {}   # slide_num → rels_xml_str
    slide_ids = {}    # slide_num → dict(media_rid, video_rid, img_rid, spid)

    with zipfile.ZipFile(str(tmp), "r") as zin:
        names = zin.namelist()
        for name in names:
            rel_m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", name)
            slide_m = re.match(r"ppt/slides/slide(\d+)\.xml", name)
            if rel_m:
                sn = int(rel_m.group(1))
                xml_str = zin.read(name).decode()
                if sn in mp3_map:
                    rids = re.findall(r'Id="rId(\d+)"', xml_str)
                    base = max((int(r) for r in rids), default=0) + 1
                    slide_ids[sn] = {
                        "media_rid": base,
                        "video_rid": base + 1,
                        "img_rid": base + 2,
                        "spid": 1000 + sn,
                    }
                slide_rels[sn] = xml_str
            elif slide_m:
                sn = int(slide_m.group(1))
                slide_xmls[sn] = zin.read(name).decode()

    MEDIA_REL = "http://schemas.microsoft.com/office/2007/relationships/media"
    VIDEO_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video"
    IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

    with zipfile.ZipFile(str(tmp), "r") as zin, \
         zipfile.ZipFile(str(out_path), "w", zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            if name == "[Content_Types].xml":
                ct = zin.read(name).decode()
                add = ""
                for sn in slide_ids:
                    part = f'/ppt/media/slide{sn}_audio.mp3'
                    if part not in ct:
                        add += f'<Override PartName="{part}" ContentType="audio/mpeg"/>'
                if 'Extension="png"' not in ct:
                    add += '<Default Extension="png" ContentType="image/png"/>'
                ct = ct.replace("</Types>", add + "</Types>")
                zout.writestr(name, ct)
                continue

            rel_m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", name)
            slide_m = re.match(r"ppt/slides/slide(\d+)\.xml", name)

            if rel_m and int(rel_m.group(1)) in slide_ids:
                sn = int(rel_m.group(1))
                ids = slide_ids[sn]
                rel_xml = slide_rels[sn]
                rels = (
                    f'<Relationship Id="rId{ids["media_rid"]}" Type="{MEDIA_REL}" '
                    f'Target="../media/slide{sn}_audio.mp3"/>'
                    f'<Relationship Id="rId{ids["video_rid"]}" Type="{VIDEO_REL}" '
                    f'Target="../media/slide{sn}_audio.mp3"/>'
                    f'<Relationship Id="rId{ids["img_rid"]}" Type="{IMAGE_REL}" '
                    f'Target="../media/slide{sn}_icon.png"/>'
                )
                rel_xml = rel_xml.replace("</Relationships>",
                                          rels + "</Relationships>")
                zout.writestr(name, rel_xml)
                zout.writestr(f"ppt/media/slide{sn}_audio.mp3", mp3_map[sn])
                zout.writestr(f"ppt/media/slide{sn}_icon.png", _ICON_PNG)
                continue

            if slide_m:
                sn = int(slide_m.group(1))
                xml_str = slide_xmls.get(sn) or zin.read(name).decode()
                adv = adv_times.get(sn, 3000)

                if sn in slide_ids:
                    ids = slide_ids[sn]
                    spid = ids["spid"]
                    pic = (
                        f'<p:pic>'
                        f'<p:nvPicPr>'
                        f'<p:cNvPr id="{spid}" name="audio{sn}.mp3">'
                        f'<a:hlinkClick r:id="" action="ppaction://media"/>'
                        f'</p:cNvPr>'
                        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
                        f'<p:nvPr>'
                        f'<a:videoFile r:link="rId{ids["video_rid"]}"/>'
                        f'<p:extLst><p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}}">'
                        f'<p14:media xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
                        f'r:embed="rId{ids["media_rid"]}"/>'
                        f'</p:ext></p:extLst>'
                        f'</p:nvPr>'
                        f'</p:nvPicPr>'
                        f'<p:blipFill><a:blip r:embed="rId{ids["img_rid"]}"/>'
                        f'<a:stretch><a:fillRect/></a:stretch></p:blipFill>'
                        f'<p:spPr><a:xfrm><a:off x="0" y="0"/>'
                        f'<a:ext cx="9144" cy="9144"/></a:xfrm>'
                        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
                        f'</p:pic>'
                    )
                    xml_str = xml_str.replace("</p:spTree>", pic + "</p:spTree>")
                    timing = _TIMING_TPL.format(spid=spid)
                else:
                    timing = ""

                transition_xml = f'<p:transition advTm="{adv}"/>'
                xml_str = xml_str.replace("</p:sld>",
                                          transition_xml + timing + "</p:sld>")
                zout.writestr(name, xml_str)
                continue

            zout.writestr(name, zin.read(name))

    tmp.unlink()


def build_title_slide(prs: Presentation, timeline: dict,
                     intro_audio_path: Path = None) -> int:
    """封面页。
    
    Args:
        prs: PPT 对象
        timeline: 时间轴数据
        intro_audio_path: 导入语音频路径（用于获取时长）
    
    Returns:
        封面页应停留的时长（毫秒），用于设置 advTm
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_GOLD

    title = timeline.get("title", "")
    author = timeline.get("author", "")
    dynasty = timeline.get("dynasty", "")

    # 标题
    _add_text_box(slide, 1, 1.8, 8, 1.2, title,
                  font_size=36, bold=True, color=WHITE)
    # 作者朝代
    if author:
        _add_text_box(slide, 1, 3.0, 8, 0.6, f"{author} · {dynasty}",
                      font_size=16, color=GOLD)

    # 底部装饰线
    line = slide.shapes.add_shape(
        1, Inches(3), Inches(3.6), Inches(4), Inches(0),
    )
    line.line.color.rgb = GOLD
    line.line.width = Pt(1)
    
    # 导入语放到备注里
    intro = timeline.get("intro", "")
    if intro:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = intro
    
    # 计算封面页停留时长 = 导入语音频时长 + 1秒余量
    intro_duration_ms = 3000  # 默认 3 秒
    if intro_audio_path and intro_audio_path.exists():
        try:
            from mutagen.mp3 import MP3
            intro_duration_ms = int(MP3(intro_audio_path).info.length * 1000)
        except Exception:
            pass
    
    return intro_duration_ms


def build_sentence_slide(prs: Presentation, sentence: dict,
                         idx: int, total: int) -> None:
    """单句内容页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WARM_CREAM

    # 备注
    narration = sentence.get("narration", "")
    if narration:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = narration

    # 左侧金色竖线装饰
    bar = slide.shapes.add_shape(
        1, Inches(0.6), Inches(0.5), Inches(0.05), Inches(3.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    # 句号
    _add_text_box(slide, 0.8, 0.3, 2, 0.4,
                  f"第 {idx} 句 / 共 {total} 句",
                  font_size=10, color=GRAY, align=PP_ALIGN.LEFT)

    # 原文
    text = sentence.get("text", "")
    _add_text_box(slide, 0.8, 0.7, 8.5, 1.0, text,
                  font_size=24, bold=True, align=PP_ALIGN.LEFT)

    # 译文
    translation = sentence.get("translation", "")
    if translation:
        _add_text_box(slide, 0.8, 1.55, 8.5, 0.7, f"译文：{translation}",
                      font_size=11, color=GRAY, align=PP_ALIGN.LEFT)

    # 关键词卡片
    keywords = sentence.get("keywords", [])
    if keywords:
        _add_text_box(slide, 0.8, 2.25, 2, 0.35, "重点词",
                      font_size=11, color=GOLD, align=PP_ALIGN.LEFT)

        card_y = 2.65
        card_w = 4.1
        card_h = 0.52
        gap_x = 0.3
        gap_y = 0.18
        cols = 2
        for i, kw in enumerate(keywords):
            col = i % cols
            row = i // cols
            cx = 0.8 + col * (card_w + gap_x)
            cy = card_y + row * (card_h + gap_y)

            color_idx = i % len(TAG_COLORS)
            _add_rounded_rect(slide, cx, cy, card_w, card_h,
                              TAG_COLORS[color_idx])

            # 词 + 释义
            _add_text_box(slide, cx + 0.15, cy, card_w - 0.3, card_h,
                          kw.get("word", ""),
                          font_size=13, bold=True, color=INK,
                          font_name="PingFang SC",
                          align=PP_ALIGN.LEFT)

            note = kw.get("note", "")
            if note:
                _add_text_box(slide, cx + 0.9, cy, card_w - 1.05, card_h,
                              note, font_size=9, color=GRAY,
                              font_name="PingFang SC",
                              align=PP_ALIGN.LEFT)


def build_pptx(timeline_path: Path, out_path: Path) -> None:
    timeline = json.loads(timeline_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    sentences = timeline.get("sentences", [])
    total = len(sentences)
    lesson_dir = timeline_path.parent

    # 封面页嵌入导入语音频
    intro_narration = timeline.get("intro", "")
    intro_audio = lesson_dir / "audio" / "00_intro.mp3"
    
    # 收集音频映射
    audio_map = {}
    if intro_narration and intro_audio.exists():
        audio_map[1] = intro_audio

    # 封面页（导入语放到备注，停留时长 = 导入语音频时长）
    intro_duration_ms = build_title_slide(prs, timeline, intro_audio)

    for i, s in enumerate(sentences, 1):
        build_sentence_slide(prs, s, i, total)
        audio_rel = s.get("audio", "")
        if audio_rel:
            audio_path = lesson_dir / audio_rel
            if audio_path.exists():
                audio_map[i + 1] = audio_path  # slide 1 = 封面, slide 2+ = 内容

    prs.save(str(out_path))

    # 后处理嵌入音频（传入封面页停留时长）
    if audio_map:
        _embed_audio_post(out_path, audio_map, intro_duration_ms)
        intro_info = "+ 1 封面导入语" if intro_narration else ""
        print(f"✓ {out_path}  ({total} 句内容 + 1 封面{intro_info}, {len(audio_map)} 段音频)")
    else:
        print(f"✓ {out_path}  ({total} 句内容 + 1 封面)")


def parse_args():
    p = argparse.ArgumentParser(description="从 timeline.json 生成讲解 PPT")
    p.add_argument("query", type=str, help="课文查询")
    p.add_argument("--lesson", type=str, default=None, help="手动指定课文名")
    p.add_argument("--out", type=str, default=None, help="输出 pptx 路径")
    return p.parse_args()


def main():
    args = parse_args()
    cfg = generate.load_config()
    
    # 意图解析（和 build_timeline.py 一致）
    intent = generate.parse_intent(args.query, cfg)
    if args.lesson:
        intent["lesson"] = args.lesson
    
    lesson_name = intent.get("lesson", "") or generate._extract_lesson(args.query)
    if args.lesson:
        lesson_name = args.lesson
    
    # 课程/章节拆分（和 build_timeline.py 一致）
    course_name, chapter_name = _split_course_chapter(lesson_name)
    course_dir = generate.lesson_name_to_pinyin(course_name)
    chapter_dir = generate.lesson_name_to_pinyin(chapter_name)
    lesson_dir = TIMELINE_DIR / course_dir / chapter_dir
    
    # 找最新的 timeline.json
    timeline_files = sorted(lesson_dir.glob("*/timeline.json"), reverse=True)
    if not timeline_files:
        print(f"✗ timeline.json 不存在: {lesson_dir}")
        print("  请先运行: python3 script/build_timeline.py '{}'".format(args.query))
        sys.exit(1)
    timeline_path = timeline_files[0]
    print(f"使用: {timeline_path}")

    out_path = Path(args.out) if args.out else timeline_path.parent / "slides.pptx"
    build_pptx(timeline_path, out_path.resolve())


if __name__ == "__main__":
    main()
