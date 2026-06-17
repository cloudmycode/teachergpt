#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""一站式生成讲解课程（合并 build_timeline.py + build_pptx.py）。

必选产出：
  data/timelines/{课程}/{章节}/script.json    讲稿（narration）
  data/timelines/{课程}/{章节}/timeline.json   时间轴（含 audio/duration，若开音频）

可选产出（通过参数开启）：
  --audio   Coze TTS 合成讲解音频，写入 timeline.json
  --pptx    从 timeline.json 生成讲解 PPT（含自动播放音频，若有）
  --mp4     仅打印用 PowerPoint 手动导出 MP4 的提示（macOS 无可靠自动转换）

流程：
  1. 复用 generate.py 的意图解析 + 课文事实 → sentences / keywords
  2. 按段落分批：每批一个段落（多句），段落间传递上下文
  3. 对每句调 DeepSeek 生成 narration → script.json（必选）
  4. （--audio）Coze TTS 合成讲解音频，线性估时对齐 → timeline.json
  5. （--pptx）每句一页：原文大字 + 关键词卡片 + 译文，封面/内容页自动播放音频
  6. （--mp4）打印手动导出提示

用法：
  python3 script/buildclass.py "诫子书 全文"
  python3 script/buildclass.py "诫子书 全文" --audio
  python3 script/buildclass.py "诫子书 全文" --audio --pptx
  python3 script/buildclass.py "诫子书 全文" --audio --pptx --mp4
  python3 script/buildclass.py "木兰词 全文" --batch-size 5 --audio
"""

import argparse
import base64
import json
import sys
import time
from pathlib import Path
from typing import Optional

SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent
sys.path.insert(0, str(PROJECT_ROOT))
sys.path.insert(0, str(SCRIPT_DIR))

import generate  # 复用意图解析、课文事实、检索、风格加载、prompt 构建

DATA_DIR = PROJECT_ROOT / "data"
TIMELINE_DIR = DATA_DIR / "timelines"


# =================================================================== TTS

def synthesize_audio(narration: str, out_path: Path) -> None:
    """Coze TTS 合成音频到 mp3。"""
    from tts.coze_tts import synthesize
    narration = narration.strip()
    if not narration:
        return
    synthesize(narration, out_path)


def read_mp3_duration(path: Path) -> float:
    """mutagen 读 MP3 时长（秒）。"""
    from mutagen.mp3 import MP3
    return MP3(path).info.length


# =================================================================== 标记解析

def parse_markers(narration: str) -> tuple[str, dict[str, int]]:
    """解析 [关键词] 标记，返回 (纯净文本, {词: 字符位置})。

    注：DeepSeek 可能在同一句里多次标记同一词（如开头引用原句时 + 讲解时）。
    这里会解析出首次标记的位置。但最终用于高亮的是按字符占比 × 音频时长的
    线性估时——该位置占 narration 长度的比例，估算音频读到该位置的时刻。
    """
    import re

    positions = {}
    clean = []
    offset = 0
    for m in re.finditer(r"\[(.+?)\]", narration):
        word = m.group(1)
        if word not in positions:
            positions[word] = m.start() - offset
        clean.append(narration[offset:m.start()])
        clean.append(word)
        offset = m.end()
    clean.append(narration[offset:])
    return "".join(clean), positions


# =================================================================== 对齐

def _narration_pos(narration: str, word: str,
                   positions: Optional[dict]) -> Optional[int]:
    """关键词在 narration 原文中的字符位置。

    优先找讲解引导词（"这个X"/"再看X"/"所谓X"），
    这是中文讲解的自然语序，比 [词] 标记更可靠。
    引导词匹配不到再退回用 [词] 标记；都没有就返回 None。
    """
    import re
    guides = r"(?:这个|再看|所谓|叫做|注意这个|先看|重点看|就是)"
    quote = r"(?:[\u201c\u201d\u2018\u2019\u300c\u300d\"'']?)"
    pattern = re.compile(f"({guides})\\s*{quote}({re.escape(word)})")
    m = pattern.search(narration)
    if m:
        return m.start(2)
    if positions and word in positions:
        return positions[word]
    return None


def align_by_ratio(narration: str, keywords: list, duration: float,
                   positions: Optional[dict] = None) -> list:
    """按字符位置线性估时定位关键词（默认对齐路径，无需 ASR）。"""
    n = len(narration)
    if n == 0 or duration <= 0:
        return []
    timings = []
    for kw in keywords:
        word = kw.get("word", "")
        if not word:
            continue
        npos = _narration_pos(narration, word, positions)
        if npos is None:
            continue
        t = npos / n * duration
        timings.append({
            "word": word,
            "note": kw.get("note", ""),
            "time": round(t, 2),
            "start": round(max(0, t - 0.5), 2),
        })
    return sorted(timings, key=lambda x: x["time"])


# =================================================================== 课程/章节拆分

def _split_course_chapter(lesson: str) -> tuple[str, str]:
    """从 lesson 名提取课程名和章节名。

    规则：
      - "德行篇第25则" → ("世说新语精读", "德行篇")
      - "言语篇" → ("世说新语精读", "言语篇")
      - "木兰诗" → ("木兰诗", "木兰诗")
      - "背影" → ("背影", "背影")
    """
    shishuo_chapters = {
        "德行篇", "言语篇", "政事篇", "文学篇", "方正篇", "雅量篇",
        "识鉴篇", "赏誉篇", "品藻篇", "规箴篇", "捷悟篇", "夙慧篇",
        "豪爽篇", "容止篇", "自新篇", "俭啬篇", "汰侈篇", "忿狷篇",
        "情礼篇", "黜免篇", "俭吝篇", "惑溺篇", "仇隙篇", "任诞篇",
        "伤逝篇", "栖逸篇", "贤媛篇", "术解篇", "巧艺篇", "知惧篇",
        "企羡篇",
    }
    for chapter in shishuo_chapters:
        if chapter in lesson:
            return ("世说新语精读", chapter)
    return (lesson, lesson)


# =================================================================== 保存

def _save_json(out_file: Path, lesson_name: str, facts: dict,
               results: list, intro_narration: str = "") -> None:
    """写入 JSON（script.json / timeline.json 同结构）。"""
    out_file.parent.mkdir(parents=True, exist_ok=True)
    obj = {
        "title": f"《{lesson_name}》精讲",
        "author": facts.get("author", ""),
        "dynasty": facts.get("dynasty", ""),
        "source": facts.get("source", ""),
        "intro": intro_narration,
        "sentences": results,
    }
    out_file.write_text(
        json.dumps(obj, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def split_sentences_by_paragraph(sentences: list, batch_size: int = 5) -> list:
    """按段落分批句子。带 paragraph 标记则按段分组，否则按 batch_size 分批。"""
    has_paragraph = any(s.get("paragraph") is not None for s in sentences)

    if has_paragraph:
        paragraphs = {}
        for s in sentences:
            para = s.get("paragraph", 0)
            paragraphs.setdefault(para, []).append(s)

        batches = []
        for para_idx in sorted(paragraphs.keys()):
            para_sentences = paragraphs[para_idx]
            if len(para_sentences) <= batch_size:
                batches.append(para_sentences)
            else:
                for i in range(0, len(para_sentences), batch_size):
                    batches.append(para_sentences[i:i + batch_size])
        return batches

    batches = []
    for i in range(0, len(sentences), batch_size):
        batches.append(sentences[i:i + batch_size])
    return batches


# =================================================================== 讲稿生成

def build_script(args, cfg, intent: dict, lesson_name: str,
                 script_file: Path):
    """生成（或读取）讲稿 script.json。

    Returns:
        (facts, results, intro_narration, total)
    """
    print(f"意图: {json.dumps(intent, ensure_ascii=False)}")

    if script_file.exists():
        print(f"讲稿已存在，跳过生成: {script_file}")
        data = json.loads(script_file.read_text(encoding="utf-8"))
        results = data.get("sentences", [])
        facts = {
            "title": data.get("title", ""),
            "author": data.get("author", ""),
            "source": data.get("source", ""),
            "dynasty": data.get("dynasty", ""),
        }
        return facts, results, data.get("intro", ""), len(results)

    # 步骤 2: 获取课文事实
    print("获取课文事实...")
    facts = generate.fetch_lesson_facts(cfg, args.query, intent)
    if not facts or not facts.get("sentences"):
        print("✗ 未获取到课文 sentences。")
        sys.exit(1)

    sentences = facts["sentences"]
    total = len(sentences)
    if args.max_sentences > 0:
        sentences = sentences[:args.max_sentences]
        total = len(sentences)
    print(
        f"课文: {facts.get('source', '')}  {facts.get('author', '')}"
        f"（{facts.get('dynasty', '')}）  共 {total} 句"
    )

    # 步骤 3: 检索相关语料
    print("检索相关语料...")
    enc = generate.Encoder()
    segments = generate.retrieve(args.query, intent, enc, rerank=args.rerank)

    # 步骤 4: 加载风格
    style = generate.load_style()

    # 步骤 5.0: 生成开头导入语
    print("\n生成开头导入语...")
    intro_system, intro_user = generate.build_intro_prompt(facts, style, segments)
    try:
        intro_narration = generate.call_deepseek(cfg, intro_system, intro_user).strip()
        print(f"  ✓ 导入语: {intro_narration[:80]}...")
    except Exception as e:
        print(f"  ✗ 导入语生成失败: {e}")
        intro_narration = ""

    # 步骤 5.1: 按段落分批生成 narration
    batch_size = args.batch_size
    batches = split_sentences_by_paragraph(sentences, batch_size)
    num_batches = len(batches)
    print(f"\n分 {num_batches} 批生成讲解（每批最多 {batch_size} 句）：")

    results = []
    prev_narration = intro_narration

    for batch_idx, batch_sentences in enumerate(batches):
        batch_start = sum(len(b) for b in batches[:batch_idx])
        batch_sentences_in_results = [
            sentences[batch_start + i] for i in range(len(batch_sentences))
        ]
        print(f"\n--- 批次 {batch_idx + 1}/{num_batches}（{len(batch_sentences)} 句）---")

        for i, s in enumerate(batch_sentences_in_results):
            idx = batch_start + i + 1
            text_preview = s.get("text", "")[:30]
            print(f"  [{idx}/{total}] {text_preview}...", end=" ", flush=True)

            t0 = time.time()
            system, user = generate.build_per_sentence_prompt(
                s, idx, total, facts, style, segments, prev_narration,
            )
            try:
                raw = generate.call_deepseek(cfg, system, user).strip()
            except Exception as e:
                print(f"✗ {e}")
                raw = ""

            narration, _ = parse_markers(raw)
            dt = time.time() - t0
            print(f"({dt:.1f}s)")

            if args.verbose and narration:
                print(f"    {narration[:150]}...")

            results.append({
                "id": idx,
                "text": s.get("text", ""),
                "translation": s.get("translation", ""),
                "keywords": s.get("keywords", []),
                "narration": narration,
            })
            prev_narration = narration

    # 先存讲稿——万一后续 TTS 崩了稿子还在
    _save_json(script_file, lesson_name, facts, results, intro_narration)
    print(f"✓ 讲稿已存: {script_file}")
    return facts, results, intro_narration, total


# =================================================================== 音频合成

def build_audio(args, results: list, intro_narration: str,
                lesson_dir: Path, total: int) -> None:
    """对讲稿逐句 TTS，回填 audio/duration 到 results（原地修改）。"""
    audio_dir = lesson_dir / "audio"
    audio_dir.mkdir(parents=True, exist_ok=True)

    # 先合成导入语音频（id=0）
    if intro_narration:
        print("\nTTS 合成导入语...")
        intro_mp3 = audio_dir / "00_intro.mp3"
        if not intro_mp3.exists():
            synthesize_audio(intro_narration, intro_mp3)
        intro_dur = read_mp3_duration(intro_mp3)
        print(f"  [0] 导入语 ({intro_dur:.1f}s)")

    print("\nTTS 合成讲解...")
    for entry in results:
        idx = entry["id"]
        narration = entry.get("narration", "")
        if not narration:
            print(f"  [{idx}/{total}] 无讲解，跳过")
            continue

        mp3_path = audio_dir / f"{idx:02d}.mp3"
        entry["audio"] = str(mp3_path.relative_to(lesson_dir))

        if mp3_path.exists():
            dur = read_mp3_duration(mp3_path)
            entry["duration"] = round(dur, 1)
            print(f"  [{idx}/{total}] 已存在，跳过 ({dur:.1f}s)")
            continue

        print(f"  [{idx}/{total}] 合成...", end=" ", flush=True)
        synthesize_audio(narration, mp3_path)
        dur = read_mp3_duration(mp3_path)
        entry["duration"] = round(dur, 1)
        print(f"{dur:.1f}s")

    # 关键词对齐（TODO: 待攻坚）
    for entry in results:
        entry["keyword_timings"] = []


# =================================================================== PPTX

# pptx 配色/常量/导入延迟到 _pptx_consts() 内，避免未装 python-pptx 时阻塞纯文本流程

# 1x1 透明 PNG，作为音频媒体对象的占位缩略图（移到画面外不显示）
_ICON_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR42mNk"
    "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)

# 自动播放音频：沿用 python-pptx/PowerPoint 原生 video media timing 结构，
# 只把开始条件从 delay="indefinite" 改为 delay="0"。
_TIMING_TPL = (
    '<p:timing><p:tnLst><p:par>'
    '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
    '<p:childTnLst><p:video><p:cMediaNode vol="80000">'
    '<p:cTn id="2" fill="hold" display="0">'
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    '</p:cTn>'
    '<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
    '</p:cMediaNode></p:video></p:childTnLst>'
    '</p:cTn></p:par></p:tnLst></p:timing>'
)


def _pptx_consts():
    """延迟导入 python-pptx，返回常用对象/常量。"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as _RGB
    consts = {
        "Inches": Inches, "Pt": Pt, "PP_ALIGN": PP_ALIGN, "RGBColor": _RGB,
        "INK": _RGB(0x3A, 0x2F, 0x1F),         # 墨色
        "WARM_CREAM": _RGB(0xF5, 0xF0, 0xE8),  # 暖米背景
        "WHITE": _RGB(0xFF, 0xFF, 0xFF),
        "GOLD": _RGB(0xC4, 0xA9, 0x7D),        # 金色
        "GRAY": _RGB(0x8C, 0x7A, 0x6B),        # 灰色
        "DARK_GOLD": _RGB(0x5C, 0x3D, 0x2E),   # 暗金色
        "SLIDE_W": Inches(10),
        "SLIDE_H": Inches(5.625),
        "TAG_COLORS": [
            _RGB(0xFA, 0xF6, 0xED),  # 米白
            _RGB(0xFB, 0xF0, 0xE0),  # 浅杏
            _RGB(0xF5, 0xEB, 0xE6),  # 浅驼
        ],
    }
    return consts


def _add_text_box(C, slide, left, top, width, height, text,
                  font_size=14, bold=False, color=None, align=None,
                  font_name="PingFang SC"):
    """添加文本框。"""
    Inches, Pt = C["Inches"], C["Pt"]
    color = color if color is not None else C["INK"]
    align = align if align is not None else C["PP_ALIGN"].CENTER
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def _add_rounded_rect(C, slide, left, top, width, height, fill_color):
    """圆角矩形。"""
    Inches = C["Inches"]
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _embed_audio_post(out_path: Path, audio_map: dict,
                      intro_duration_ms: int = 3000) -> None:
    """后处理 pptx：注入 PowerPoint 原生媒体对象，并设为自动播放。"""
    import zipfile
    import shutil
    import re
    from mutagen.mp3 import MP3

    tmp = out_path.with_suffix(".tmp")
    shutil.move(str(out_path), str(tmp))

    adv_times = {}
    mp3_map = {}
    for sn, audio_path in audio_map.items():
        if audio_path.exists():
            mp3_map[sn] = audio_path.read_bytes()
            adv_times[sn] = int(MP3(audio_path).info.length * 1000)

    if 1 not in adv_times:
        adv_times[1] = intro_duration_ms
    # intro 读完停 1 秒再翻页
    if 1 in adv_times and adv_times[1] > 0:
        adv_times[1] += 1000

    slide_xmls = {}
    slide_rels = {}
    slide_ids = {}

    with zipfile.ZipFile(str(tmp), "r") as zin:
        names = zin.namelist()
        for name in names:
            rel_m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", name)
            slide_m = re.match(r"ppt/slides/slide(\d+)\.xml", name)
            if rel_m:
                sn = int(rel_m.group(1))
                xml_str = zin.read(name).decode()
                if sn in mp3_map:
                    rids = re.findall(r'Id="rId(\d+)"', xml_str)
                    base = max((int(r) for r in rids), default=0) + 1
                    slide_ids[sn] = {
                        "media_rid": base,
                        "video_rid": base + 1,
                        "img_rid": base + 2,
                        "spid": 1000 + sn,
                    }
                slide_rels[sn] = xml_str
            elif slide_m:
                sn = int(slide_m.group(1))
                slide_xmls[sn] = zin.read(name).decode()

    MEDIA_REL = "http://schemas.microsoft.com/office/2007/relationships/media"
    VIDEO_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video"
    IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

    with zipfile.ZipFile(str(tmp), "r") as zin, \
         zipfile.ZipFile(str(out_path), "w", zipfile.ZIP_DEFLATED) as zout:
        for name in names:
            if name == "[Content_Types].xml":
                ct = zin.read(name).decode()
                add = ""
                for sn in slide_ids:
                    part = f'/ppt/media/slide{sn}_audio.mp3'
                    if part not in ct:
                        add += f'<Override PartName="{part}" ContentType="audio/mpeg"/>'
                if 'Extension="png"' not in ct:
                    add += '<Default Extension="png" ContentType="image/png"/>'
                ct = ct.replace("</Types>", add + "</Types>")
                zout.writestr(name, ct)
                continue

            rel_m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", name)
            slide_m = re.match(r"ppt/slides/slide(\d+)\.xml", name)

            if rel_m and int(rel_m.group(1)) in slide_ids:
                sn = int(rel_m.group(1))
                ids = slide_ids[sn]
                rel_xml = slide_rels[sn]
                rels = (
                    f'<Relationship Id="rId{ids["media_rid"]}" Type="{MEDIA_REL}" '
                    f'Target="../media/slide{sn}_audio.mp3"/>'
                    f'<Relationship Id="rId{ids["video_rid"]}" Type="{VIDEO_REL}" '
                    f'Target="../media/slide{sn}_audio.mp3"/>'
                    f'<Relationship Id="rId{ids["img_rid"]}" Type="{IMAGE_REL}" '
                    f'Target="../media/slide{sn}_icon.png"/>'
                )
                rel_xml = rel_xml.replace("</Relationships>",
                                          rels + "</Relationships>")
                zout.writestr(name, rel_xml)
                zout.writestr(f"ppt/media/slide{sn}_audio.mp3", mp3_map[sn])
                zout.writestr(f"ppt/media/slide{sn}_icon.png", _ICON_PNG)
                continue

            if slide_m:
                sn = int(slide_m.group(1))
                xml_str = slide_xmls.get(sn) or zin.read(name).decode()
                adv = adv_times.get(sn, 3000)

                if sn in slide_ids:
                    ids = slide_ids[sn]
                    spid = ids["spid"]
                    pic = (
                        f'<p:pic>'
                        f'<p:nvPicPr>'
                        f'<p:cNvPr id="{spid}" name="audio{sn}.mp3">'
                        f'<a:hlinkClick r:id="" action="ppaction://media"/>'
                        f'</p:cNvPr>'
                        f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
                        f'<p:nvPr>'
                        f'<a:videoFile r:link="rId{ids["video_rid"]}"/>'
                        f'<p:extLst><p:ext uri="{{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}}">'
                        f'<p14:media xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
                        f'r:embed="rId{ids["media_rid"]}"/>'
                        f'</p:ext></p:extLst>'
                        f'</p:nvPr>'
                        f'</p:nvPicPr>'
                        f'<p:blipFill><a:blip r:embed="rId{ids["img_rid"]}"/>'
                        f'<a:stretch><a:fillRect/></a:stretch></p:blipFill>'
                        f'<p:spPr><a:xfrm><a:off x="0" y="0"/>'
                        f'<a:ext cx="9144" cy="9144"/></a:xfrm>'
                        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
                        f'</p:pic>'
                    )
                    xml_str = xml_str.replace("</p:spTree>", pic + "</p:spTree>")
                    timing = _TIMING_TPL.format(spid=spid)
                else:
                    timing = ""

                transition_xml = f'<p:transition advTm="{adv}"/>'
                xml_str = xml_str.replace("</p:sld>",
                                          transition_xml + timing + "</p:sld>")
                zout.writestr(name, xml_str)
                continue

            zout.writestr(name, zin.read(name))

    tmp.unlink()


def build_title_slide(C, prs, timeline: dict, intro_audio_path: Path = None) -> int:
    """封面页。返回封面页应停留的时长（毫秒）。"""
    Inches, Pt = C["Inches"], C["Pt"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["DARK_GOLD"]

    title = timeline.get("title", "")
    author = timeline.get("author", "")
    dynasty = timeline.get("dynasty", "")

    _add_text_box(C, slide, 1, 1.8, 8, 1.2, title,
                  font_size=36, bold=True, color=C["WHITE"])
    if author:
        _add_text_box(C, slide, 1, 3.0, 8, 0.6, f"{author} · {dynasty}",
                      font_size=16, color=C["GOLD"])

    line = slide.shapes.add_shape(
        1, Inches(3), Inches(3.6), Inches(4), Inches(0),
    )
    line.line.color.rgb = C["GOLD"]
    line.line.width = Pt(1)

    intro = timeline.get("intro", "")
    if intro:
        slide.notes_slide.notes_text_frame.text = intro

    _add_text_box(C, slide, 0, 5.1, 10, 0.35,
                  "本课件由精讲课AI辅助生成，了解生成原理：www.jingjiangke.com",
                  font_size=8, color=C["GOLD"])

    intro_duration_ms = 3000
    if intro_audio_path and intro_audio_path.exists():
        try:
            from mutagen.mp3 import MP3
            intro_duration_ms = int(MP3(intro_audio_path).info.length * 1000)
        except Exception:
            pass
    return intro_duration_ms


def build_sentence_slide(C, prs, sentence: dict, idx: int, total: int) -> None:
    """单句内容页。"""
    PP_ALIGN = C["PP_ALIGN"]
    Inches = C["Inches"]
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["WARM_CREAM"]

    narration = sentence.get("narration", "")
    if narration:
        slide.notes_slide.notes_text_frame.text = narration

    bar = slide.shapes.add_shape(
        1, Inches(0.6), Inches(0.5), Inches(0.05), Inches(3.2),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C["GOLD"]
    bar.line.fill.background()

    _add_text_box(C, slide, 0.8, 0.3, 2, 0.4,
                  f"第 {idx} 句 / 共 {total} 句",
                  font_size=10, color=C["GRAY"], align=PP_ALIGN.LEFT)

    text = sentence.get("text", "")
    _add_text_box(C, slide, 0.8, 0.7, 8.5, 1.0, text,
                  font_size=24, bold=True, align=PP_ALIGN.LEFT)

    translation = sentence.get("translation", "")
    if translation:
        _add_text_box(C, slide, 0.8, 1.55, 8.5, 0.7, f"译文：{translation}",
                      font_size=11, color=C["GRAY"], align=PP_ALIGN.LEFT)

    _add_text_box(C, slide, 0, 5.1, 10, 0.35,
                  "本课件由精讲课AI辅助生成，了解生成原理：www.jingjiangke.com",
                  font_size=8, color=C["GRAY"])

    keywords = sentence.get("keywords", [])
    if keywords:
        _add_text_box(C, slide, 0.8, 2.25, 2, 0.35, "重点词",
                      font_size=11, color=C["GOLD"], align=PP_ALIGN.LEFT)

        card_y, card_w, card_h = 2.65, 4.1, 0.52
        gap_x, gap_y, cols = 0.3, 0.18, 2
        for i, kw in enumerate(keywords):
            col, row = i % cols, i // cols
            cx = 0.8 + col * (card_w + gap_x)
            cy = card_y + row * (card_h + gap_y)

            _add_rounded_rect(C, slide, cx, cy, card_w, card_h,
                              C["TAG_COLORS"][i % len(C["TAG_COLORS"])])
            _add_text_box(C, slide, cx + 0.15, cy, card_w - 0.3, card_h,
                          kw.get("word", ""),
                          font_size=13, bold=True, color=C["INK"],
                          align=PP_ALIGN.LEFT)
            note = kw.get("note", "")
            if note:
                _add_text_box(C, slide, cx + 0.9, cy, card_w - 1.05, card_h,
                              note, font_size=9, color=C["GRAY"],
                              align=PP_ALIGN.LEFT)


def build_pptx(timeline_path: Path, out_path: Path) -> None:
    """从 timeline.json 生成讲解 PPT（含自动播放音频，若有）。"""
    from pptx import Presentation
    C = _pptx_consts()

    timeline = json.loads(timeline_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = C["SLIDE_W"]
    prs.slide_height = C["SLIDE_H"]

    sentences = timeline.get("sentences", [])
    total = len(sentences)
    lesson_dir = timeline_path.parent

    intro_narration = timeline.get("intro", "")
    intro_audio = lesson_dir / "audio" / "00_intro.mp3"

    audio_map = {}
    if intro_narration and intro_audio.exists():
        audio_map[1] = intro_audio

    intro_duration_ms = build_title_slide(C, prs, timeline, intro_audio)

    for i, s in enumerate(sentences, 1):
        build_sentence_slide(C, prs, s, i, total)
        audio_rel = s.get("audio", "")
        if audio_rel:
            audio_path = lesson_dir / audio_rel
            if audio_path.exists():
                audio_map[i + 1] = audio_path  # slide 1 = 封面

    prs.save(str(out_path))

    if audio_map:
        _embed_audio_post(out_path, audio_map, intro_duration_ms)
        intro_info = "+ 1 封面导入语" if intro_narration else ""
        print(f"✓ {out_path}  ({total} 句内容 + 1 封面{intro_info}, {len(audio_map)} 段音频)")
    else:
        print(f"✓ {out_path}  ({total} 句内容 + 1 封面)")


# =================================================================== MP4 提示

def print_mp4_hint(pptx_path: Path) -> None:
    """macOS 上 pptx→mp4 无可靠自动方案，打印手动导出步骤。"""
    print("\n--- 导出 MP4（手动）---")
    if pptx_path.exists():
        print(f"PPTX: {pptx_path}")
    else:
        print(f"✗ 未找到 PPTX（请加 --pptx 先生成）: {pptx_path}")
        return
    print("macOS PowerPoint 导出步骤：")
    print("  1. 用 PowerPoint 打开上面的 pptx")
    print("  2. 菜单：文件 → 导出")
    print("  3. 文件格式选 MP4，设置分辨率")
    print("  4. 导出（每页已设自动切换/嵌入音频，会按音频时长推进）")


# =================================================================== main

def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="一站式生成讲解课程：script/timeline（必选），可选音频/PPTX/MP4 提示"
    )
    p.add_argument("query", type=str, help="课文查询，如'诫子书 全文'")
    p.add_argument("--lesson", type=str, default=None, help="手动指定课文名")
    p.add_argument("--max-sentences", type=int, default=0,
                   help="最多处理前 N 句（0=全量），调试用")
    p.add_argument("--batch-size", type=int, default=5,
                   help="每批最大句数（默认 5）")
    p.add_argument("--rerank", action="store_true", help="启用 reranker 精排")
    p.add_argument("--verbose", "-v", action="store_true", help="打印讲解内容摘要")
    # 可选产出
    p.add_argument("--audio", action="store_true",
                   help="合成讲解音频并写入 timeline.json")
    p.add_argument("--pptx", action="store_true",
                   help="从 timeline.json 生成讲解 PPT")
    p.add_argument("--pptx-out", type=str, default=None,
                   help="PPT 输出路径（默认章节目录下 slides.pptx）")
    p.add_argument("--mp4", action="store_true",
                   help="打印用 PowerPoint 手动导出 MP4 的提示")
    return p.parse_args()


def main() -> None:
    args = parse_args()

    cfg = generate.load_config()
    if not cfg["api_key"]:
        print("✗ 未找到 DeepSeek api_key。")
        sys.exit(1)

    # ---- 意图解析（仅一次）+ 输出目录 ----
    intent = generate.parse_intent(args.query, cfg)
    intent["intent"] = "精读讲解"
    if args.lesson:
        intent["lesson"] = args.lesson
    lesson_name = (args.lesson or intent.get("lesson", "")
                   or generate._extract_lesson(args.query))
    course_name, chapter_name = _split_course_chapter(lesson_name)
    course_dir = generate.lesson_name_to_pinyin(course_name)
    chapter_dir = generate.lesson_name_to_pinyin(chapter_name)
    lesson_dir = TIMELINE_DIR / course_dir / chapter_dir
    script_file = lesson_dir / "script.json"
    timeline_file = lesson_dir / "timeline.json"

    # ---- 步骤 1（必选）：讲稿 ----
    facts, results, intro_narration, total = build_script(
        args, cfg, intent, lesson_name, script_file
    )

    # ---- 步骤 2（可选）：音频 ----
    if args.audio:
        build_audio(args, results, intro_narration, lesson_dir, total)

    # ---- 步骤 3（必选）：timeline.json ----
    _save_json(timeline_file, lesson_name, facts, results, intro_narration)
    print(f"✓ timeline: {timeline_file}")

    ok = sum(1 for r in results if r.get("narration"))
    chars = sum(len(r.get("narration", "")) for r in results)
    dur = sum(r.get("duration", 0) for r in results)
    print(f"  有效讲解: {ok}/{len(results)} 句  {chars} 字  {dur:.0f}s")

    # ---- 步骤 4（可选）：PPTX ----
    pptx_path = (Path(args.pptx_out).resolve() if args.pptx_out
                 else lesson_dir / "slides.pptx")
    if args.pptx:
        print(f"\n生成 PPT: {timeline_file}")
        build_pptx(timeline_file, pptx_path)

    # ---- 步骤 5（可选）：MP4 提示 ----
    if args.mp4:
        print_mp4_hint(pptx_path)


if __name__ == "__main__":
    main()
